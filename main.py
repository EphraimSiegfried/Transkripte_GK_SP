import os
import requests
import time
import pathlib
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from zipfile import ZipFile
import re

TRANSCRIPT_PPTX_DIR = ORIGINAL_PPTX_DIR = "Transcript"
TRANSCRIPT = os.path.join(TRANSCRIPT_PPTX_DIR, "transcript.txt")


def read_audio_file(file_path, chunk_size=5242880):
    with open(file_path, "rb") as f:
        while True:
            data = f.read(chunk_size)
            if not data:
                break
            yield data


# Upload file to AssemblyAI
def transcribe(file_path):
    # Upload file to AssemblyAI
    upload_endpoint = "https://api.assemblyai.com/v2/upload"
    headers = {"authorization": os.environ["ASSEMBLYAI_API_KEY"]}
    upload_response = requests.post(upload_endpoint, headers=headers, data=read_audio_file(file_path))
    audio_url = upload_response.json()["upload_url"]

    # Start transcription
    transcription_endpoint = "https://api.assemblyai.com/v2/transcript"
    json = {"audio_url": audio_url, "language_code": "de"}
    transcription_response = requests.post(transcription_endpoint, json=json, headers=headers)
    transcription_id = transcription_response.json()["id"]
    polling_endpoint = f"https://api.assemblyai.com/v2/transcript/{transcription_id}"

    # Poll for transcription
    while True:
        response = requests.get(polling_endpoint, headers=headers)
        status = response.json()["status"]
        if status == "completed":
            transcription = response.json()["text"]
            break
        elif status == "error":
            raise Exception("Transcription failed")
        else:
            print(status)
            time.sleep(1)
    return transcription


def write_text_in_notes_section(text, slide):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def write_transcription(pptx_file):
    """
    Iterates through the powerpoint and if it
    finds a audio, it will transcribe it and write the transcription in the notes section"""
    prs = Presentation(pptx_file)
    with tempfile.TemporaryDirectory() as tmp_dir:
        audio_files = list_audiofiles(pptx_file, tmp_dir)
        n = 0
        for slide in prs.slides:
            n += 1
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    print(f"Slide {n} contains audio")
                    audio_file = audio_files.pop(0)
                    print("Begin transcription of ", audio_file)
                    text = structure_text(transcribe(audio_file))
                    write_text_in_notes_section(text, slide)
                    print("Transcription has been written to notes section\n")
                    append_to_txt_file(TRANSCRIPT, f"Slide {n}:\n{text}")
                    break
    prs.save(pptx_file)


def append_to_txt_file(file, text):
    """Appends text with a new line to a text file"""
    with open(file, "a") as f:
        f.write(text + "\n")


def list_audiofiles(pptx_file, tmp_dir):
    with ZipFile(pptx_file) as pptx:
        media_dir = pathlib.Path("ppt/media")
        pptx.extractall(tmp_dir,
                        members=[f for f in pptx.namelist() if f.startswith(str(media_dir)) and f.endswith(".m4a")])
        audio_files = [str(f) for f in (tmp_dir / media_dir).iterdir() if f.is_file()]
        # Sort audio_files by the number after the string "media" and before the string ".m4a"
        audio_files.sort(key=lambda f: int(re.search(r"media(\d+).m4a", f).group(1)))
        return audio_files


def structure_text(text):
    # Make two new lines every six sentences
    sentences = re.split(r"(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s", text)
    return "\n\n".join([" ".join(sentences[i:i + 6]) for i in range(0, len(sentences), 6)])


def text_to_markdown(txt_file):
    """ Converts a txt file into a markdown. Whenever a line starts with Powerpoint: it will be a header. Whenever
    a line starts with Slide it will be a subheader. It should create a new file """
    with open(txt_file, "r") as f:
        text = f.read()
    text = re.sub(r"Powerpoint: (.*)\n", r"# \1\n", text)
    text = re.sub(r"Slide (\d+):\n", r"## Slide \1\n", text)
    with open("transcript.md", "w") as f:
        f.write(text)


def main():
    for file in sorted(os.listdir(ORIGINAL_PPTX_DIR)):
        print("Transcribing ", file)
        append_to_txt_file(TRANSCRIPT, f"\nPowerpoint: {file}\n")
        write_transcription(os.path.join(TRANSCRIPT_PPTX_DIR, file))
        print("Transcription complete of ", file)


if __name__ == "__main__":
    main()
